# ============================================================
#  YADOKARI Trailer House Proposal  ─  v3 (Brand Redesign)
#  Design: White / #111 / #666 / #DDD  — gothic only
# ============================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── palette ─────────────────────────────────────────────────
W   = RGBColor(0xFF, 0xFF, 0xFF)   # white
INK = RGBColor(0x11, 0x11, 0x11)   # near-black
G1  = RGBColor(0x44, 0x44, 0x44)   # dark gray  (body)
G2  = RGBColor(0x88, 0x88, 0x88)   # mid gray   (secondary)
G3  = RGBColor(0xBB, 0xBB, 0xBB)   # light gray (captions)
DIV = RGBColor(0xDD, 0xDD, 0xDD)   # divider
BG  = RGBColor(0xF6, 0xF5, 0xF3)   # warm off-white bg
BLK = RGBColor(0x00, 0x00, 0x00)   # pure black

FONT = "Yu Gothic"   # gothic / sans-serif
# ── helpers ─────────────────────────────────────────────────
def S(dark=False):
    sl = prs.slides.add_slide(BLANK)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = INK if dark else W
    return sl

def R(sl, l, t, w, h, fc=DIV, lc=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def L(sl, x1, y1, x2, y2, color=DIV, pt=0.75):
    c = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(pt)

def T(sl, text, l, t, w, h, sz=13, bold=False, color=INK,
      align=PP_ALIGN.LEFT, wrap=True, italic=False):
    b = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = FONT
    return b

def MT(sl, l, t, w, h, wrap=True):
    b = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = wrap
    return tf

def AP(tf, text, sz=13, bold=False, color=INK, align=PP_ALIGN.LEFT,
       sb=0, sa=0, italic=False):
    p = tf.paragraphs[0] if (len(tf.paragraphs)==1 and tf.paragraphs[0].text=='') else tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = FONT
    return p

def HBAR(sl, title, subtitle=None):
    """Page header with thick black left-edge stripe."""
    R(sl, 0.0, 0.0, 0.08, 7.5, fc=INK)          # left stripe
    T(sl, title,    0.25, 0.25, 9.0, 0.75, sz=30, bold=True)
    if subtitle:
        T(sl, subtitle, 0.25, 0.95, 9.0, 0.38, sz=13, color=G2)
    L(sl, 0.25, 1.42 if subtitle else 1.15, 13.0, 1.42 if subtitle else 1.15, color=DIV)

def PHOTO(sl, l, t, w, h, label="写真エリア"):
    R(sl, l, t, w, h, fc=RGBColor(0xE2,0xE2,0xE0))
    T(sl, f"[ {label} ]", l, t + h/2 - 0.2, w, 0.4,
      sz=10, color=G2, align=PP_ALIGN.CENTER)

def stat_block(sl, x, y, number, unit, label, sz_num=54):
    """Large-number stat block."""
    T(sl, number, x, y,       3.0, 1.0, sz=sz_num, bold=True, color=INK, align=PP_ALIGN.CENTER)
    T(sl, unit,   x, y+0.95,  3.0, 0.35, sz=13, color=G2, align=PP_ALIGN.CENTER)
    T(sl, label,  x, y+1.25,  3.0, 0.35, sz=11, color=G2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# 01 TITLE
# ═══════════════════════════════════════════════════════════════════
sl = S(dark=True)
# White block right 60%
R(sl, 5.3, 0.0, 8.03, 7.5, fc=W)
# Left: bold headline
T(sl, "小さいほど、", 0.45, 1.4, 4.7, 1.2, sz=46, bold=True, color=W)
T(sl, "豊かだ。",     0.45, 2.6, 4.7, 1.2, sz=46, bold=True, color=W)
L(sl, 0.45, 3.85, 4.8, 3.85, color=G2, pt=0.5)
T(sl, "YADOKARI Trailer House Proposal", 0.45, 4.05, 4.7, 0.45,
  sz=13, color=G2)
T(sl, "YADOKARI株式会社", 0.45, 6.9, 4.7, 0.35, sz=11, color=G2)
# Right: photo
PHOTO(sl, 5.6, 0.4, 7.4, 6.6, "表紙写真 — トレーラーハウス外観")
T(sl, "yadokari.net", 11.5, 7.1, 1.6, 0.3, sz=10, color=G2)

# ═══════════════════════════════════════════════════════════════════
# 02 MISSION + STATS
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "ミッション", "世界を変える、暮らしを創る。")
# Big quote
T(sl, "「世界を変える、暮らしを創る」の実現を目指す", 0.5, 1.7, 12.3, 0.55, sz=17, bold=True)
T(sl, "ソーシャルデザインカンパニー。小さな空間の中に広がる、大きな可能性。",
  0.5, 2.25, 12.3, 0.55, sz=17, bold=True)
T(sl, "多様な生き方・暮らしのカタチに触れ、未来のライフスタイルを探求し提供し続けます。",
  0.5, 2.9, 12.3, 0.45, sz=13, color=G1)
L(sl, 0.5, 3.55, 12.8, 3.55, color=DIV)
# Stats — large numbers
T(sl, "実　績", 0.5, 3.75, 2.5, 0.4, sz=11, bold=True, color=G2)
stat_block(sl,  0.3, 4.1,  "●●",   "台以上",    "トレーラーハウス\n販売実績",    sz_num=52)
stat_block(sl,  3.5, 4.1,  "●●",   "件以上",    "導入施設",                      sz_num=52)
stat_block(sl,  6.7, 4.1,  "●●",   "媒体",      "メディア掲載\nNHK他",          sz_num=52)
stat_block(sl, 10.0, 4.1,  "2012", "年設立",    "神奈川県横浜市",               sz_num=44)

# ═══════════════════════════════════════════════════════════════════
# 03 WHY NOW — market background
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "なぜ今、トレーラーハウスなのか？", "市場背景")

cols = [
    ("インバウンド\n回復",
     "訪日観光客の急増に伴い、グランピング・体験型宿泊の需要が急拡大。"),
    ("移住・\n二拠点生活",
     "地方移住や二拠点生活トレンドが加速。手軽に始められる住まいのニーズが増大。"),
    ("遊休地の\n有効活用",
     "空き地・農地・山林などの遊休資産を収益化したいオーナーが増加。"),
    ("建設業の\n課題",
     "職人不足・工期長期化・コスト高騰が深刻化。スピードと低コストが求められる。"),
]
for i, (ttl, body) in enumerate(cols):
    x = 0.45 + i * 3.23
    R(sl, x, 1.6, 3.0, 1.3, fc=INK)
    T(sl, ttl, x+0.1, 1.65, 2.8, 1.2, sz=14, bold=True, color=W, align=PP_ALIGN.CENTER)
    T(sl, body, x+0.1, 3.05, 2.8, 1.5, sz=12, color=G1, wrap=True)

L(sl, 0.45, 4.75, 12.85, 4.75, color=DIV)
T(sl, "YADOKARIのトレーラーハウスは、これらすべての課題に応える「可動産」という新しい選択肢です。",
  0.5, 4.92, 12.3, 0.5, sz=16, bold=True)
tf = MT(sl, 0.5, 5.55, 12.3, 1.5)
AP(tf, "・初期費用を抑えながら高品質な空間を短期間で実現　　・移動できるから、土地を選ばない柔軟な事業展開が可能", sz=12, color=G1)
AP(tf, "・建築物でない「車両」扱いで、法規制・税制面のハードルを低減　　・中古売却（出口戦略）まで一貫サポート", sz=12, color=G1, sb=3)

# ═══════════════════════════════════════════════════════════════════
# 04 MOVABLE PROPERTY — merits + comparison table
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "可動産という選択肢")

merits = [
    ("機動力",     "基礎不要。場所を選ばない\n柔軟な運用。"),
    ("コスト",     "従来の建築物に比べ\n初期費用を大幅に抑制。"),
    ("スピード",   "製造〜納品まで\n最短3〜4ヶ月。"),
    ("法制メリット","条件次第で「車両」扱い。\n税制・手続のハードルを低減。"),
]
for i, (ttl, body) in enumerate(merits):
    x = 0.45 + i * 3.23
    if i: L(sl, x-0.12, 1.5, x-0.12, 3.8, color=DIV)
    T(sl, f"0{i+1}", x+0.05, 1.55, 0.7, 0.55, sz=24, bold=True, color=DIV)
    T(sl, ttl,  x+0.05, 2.1,  2.9, 0.5, sz=16, bold=True)
    T(sl, body, x+0.05, 2.65, 2.9, 0.9, sz=12, color=G1, wrap=True)

L(sl, 0.45, 3.9, 12.85, 3.9)
T(sl, "従来の建築物との比較", 0.5, 4.05, 4.0, 0.35, sz=11, bold=True, color=G2)

cols_w = [3.0, 3.5, 3.5]
cols_x = [0.5, 3.6, 7.2]
headers = ["", "従来の建築物", "トレーラーハウス"]
tbl_rows = [
    ["初期費用",  "数千万〜数億円",    "749万円〜"],
    ["工  期",    "6ヶ月〜1年以上",   "約3〜4ヶ月"],
    ["移設・撤去","ほぼ不可",          "可能（再設置・売却）"],
    ["建築確認",  "必要",              "不要（条件付き）※"],
    ["固定資産税","課税",              "非課税（条件付き）※"],
    ["出口戦略",  "解体コスト発生",    "中古売却が可能"],
]
RH = 0.35
YS = 4.48
for ci, (hd, cx, cw) in enumerate(zip(headers, cols_x, cols_w)):
    R(sl, cx, YS, cw-0.05, RH, fc=INK if ci else RGBColor(0xF0,0xF0,0xEE))
    T(sl, hd, cx+0.08, YS+0.06, cw-0.15, RH-0.08, sz=11, bold=True,
      color=W if ci else G2, align=PP_ALIGN.CENTER)
for ri, row in enumerate(tbl_rows):
    y = YS + RH + ri*RH
    for ci, (cell, cx, cw) in enumerate(zip(row, cols_x, cols_w)):
        fc = (RGBColor(0xF0,0xF6,0xF0) if ci==2 else
              (W if ri%2 else RGBColor(0xFA,0xFA,0xF8)))
        R(sl, cx, y, cw-0.05, RH-0.02, fc=fc)
        T(sl, cell, cx+0.1, y+0.06, cw-0.18, RH-0.1, sz=11, bold=(ci==2),
          align=PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)
T(sl, "※一定条件下において。詳細はご相談ください。", 0.5, 7.2, 8.0, 0.25, sz=9, color=G3)

# ═══════════════════════════════════════════════════════════════════
# 05 USE CASES
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "多様な活用シーン")

scenes = [
    ("宿泊施設・グランピング",
     "高いデザイン性で宿泊単価を向上。複数台導入によるヴィレッジ展開にも最適。"),
    ("店舗・カフェ",
     "洗練された外観で集客力を強化。スモールスタートでの新規事業に。"),
    ("オフィス・サテライト",
     "節税対策や、自然の中でのサテライトオフィス構築に。"),
    ("住まい・セカンドハウス",
     "趣味の空間や、二拠点居住のベースキャンプとして。"),
]
positions = [(0.45, 1.55), (6.9, 1.55), (0.45, 4.4), (6.9, 4.4)]
for (x, y), (ttl, body) in zip(positions, scenes):
    PHOTO(sl, x, y, 5.9, 1.75)
    R(sl, x, y+1.75, 5.9, 0.08, fc=INK)    # accent underline
    T(sl, ttl,  x, y+1.9,  5.9, 0.45, sz=14, bold=True)
    T(sl, body, x, y+2.4,  5.9, 0.65, sz=12, color=G1, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# 06 YADOKARI STRENGTHS
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "YADOKARIの強み")

strengths = [
    ("洗練された\nデザイン",
     "空間の豊かさを追求したデザイン性。\n用途に合わせたセミオーダー提案。"),
    ("圧倒的な\nコストパフォーマンス",
     "無駄を省いた設計と標準化で\n初期投資を最小化し、高ROIを実現。"),
    ("一気通貫の\nエコシステム",
     "企画・設計から納品、中古売買\nプラットフォームまで完全サポート。"),
]
for i, (ttl, body) in enumerate(strengths):
    x = 0.45 + i * 4.35
    if i: L(sl, x-0.12, 1.55, x-0.12, 7.1, color=DIV)
    R(sl, x, 1.55, 4.1, 0.06, fc=INK)    # top accent bar
    T(sl, f"0{i+1}", x+0.1, 1.75, 3.8, 0.7, sz=40, bold=True, color=DIV)
    T(sl, ttl,  x+0.1, 2.45, 3.8, 1.0, sz=22, bold=True, align=PP_ALIGN.LEFT)
    T(sl, body, x+0.1, 3.65, 3.8, 1.8, sz=13, color=G1, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# 07 CASE STUDIES (new)
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "導入事例", "実際の活用シーンと成果")

cases = [
    {
        "title": "事例①　グランピング施設（地方リゾート）",
        "model": "ROADIE × 5台",
        "issue": "課題：短期間で宿泊棟を増やしたかった",
        "kpi":   [("工 期", "4ヶ月", "従来比 1/3"), ("稼働率", "80%", "導入後6ヶ月"), ("宿泊単価", "×1.5", "従来比")],
        "voice": "「建築物では実現できなかったスピードとデザイン性を\n両立できました。」── ○○施設 オーナー様",
    },
    {
        "title": "事例②　サテライトオフィス（首都圏近郊）",
        "model": "Tinys INSPIRATION 9m × 2台",
        "issue": "課題：テレワーク対応の執務スペースが必要",
        "kpi":   [("設置完了", "2ヶ月", "工期"), ("節税効果", "あり", "固定資産税"), ("社員満足", "UP", "採用PRにも活用")],
        "voice": "「こんなに早く、コストを抑えて快適な空間が\n手に入るとは思いませんでした。」── ○○株式会社",
    },
]
for i, c in enumerate(cases):
    x = 0.45 + i * 6.45
    R(sl, x, 1.58, 6.1, 0.06, fc=INK)
    T(sl, c["title"], x, 1.75, 6.0, 0.4, sz=13, bold=True)
    T(sl, c["model"], x, 2.2,  6.0, 0.35, sz=11, color=G2)
    T(sl, c["issue"], x, 2.65, 6.0, 0.35, sz=12, bold=True, color=G1)
    # KPI numbers
    for j, (lbl, num, sub) in enumerate(c["kpi"]):
        kx = x + j*2.0
        T(sl, num, kx, 3.1, 1.85, 0.75, sz=34, bold=True, align=PP_ALIGN.CENTER)
        T(sl, lbl, kx, 3.85, 1.85, 0.3, sz=10, color=G2, align=PP_ALIGN.CENTER)
        T(sl, sub, kx, 4.15, 1.85, 0.3, sz=10, color=G3, align=PP_ALIGN.CENTER)
    L(sl, x, 4.55, x+5.9, 4.55, color=DIV)
    T(sl, c["voice"], x, 4.65, 5.9, 0.95, sz=11, color=G1, italic=True, wrap=True)

T(sl, "※事例写真・詳細データは別途ご用意できます。実績数値はすべてモデルケースです。",
  0.45, 7.2, 12.4, 0.25, sz=9, color=G3)

# ═══════════════════════════════════════════════════════════════════
# 08 PRODUCT LINEUP overview
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "プロダクトラインナップ")

models = [
    ("Tinys INSPIRATION", "定番ベーシックモデル",
     "低価格・高品質。\nサイズバリエーションが豊富。",   "7,490,000"),
    ("ROADIE",            "自然と繋がるモデル",
     "周囲の自然とシームレスに繋がる\n開放的デザイン。",  "10,250,000"),
    ("MIGRA",             "上質な木造シリーズ",
     "木の温かみと開放的リビング。\nプレミアム体験を提供。","9,950,000"),
]
for i, (name, sub, desc, price) in enumerate(models):
    x = 0.45 + i * 4.35
    if i: L(sl, x-0.12, 1.55, x-0.12, 7.2, color=DIV)
    PHOTO(sl, x, 1.55, 4.1, 2.6)
    R(sl, x, 4.15, 4.1, 0.06, fc=INK)
    T(sl, name, x, 4.28, 4.0, 0.5, sz=18, bold=True)
    T(sl, sub,  x, 4.78, 4.0, 0.35, sz=12, color=G2)
    T(sl, desc, x, 5.18, 4.0, 0.85, sz=12, color=G1, wrap=True)
    T(sl, "¥", x,     6.15, 0.45, 0.6, sz=14, bold=True, color=G2)
    T(sl, price, x+0.4, 6.0, 3.5, 0.7, sz=26, bold=True)
    T(sl, "円〜（税抜）", x+0.4, 6.7, 3.5, 0.3, sz=11, color=G2)

# ═══════════════════════════════════════════════════════════════════
# 09 Tinys INSPIRATION detail
# ═══════════════════════════════════════════════════════════════════
sl = S()
PHOTO(sl, 0.0, 0.0, 5.8, 7.5, "Tinys INSPIRATION")
R(sl, 5.8, 0.0, 7.53, 7.5, fc=W)  # right white panel
R(sl, 5.8, 0.0, 0.06, 7.5, fc=INK)  # border stripe
T(sl, "Tinys INSPIRATION", 6.1, 0.45, 7.0, 0.8, sz=28, bold=True)
T(sl, "定番ベーシックモデル", 6.1, 1.3, 7.0, 0.38, sz=14, color=G2)
L(sl, 6.1, 1.8, 13.1, 1.8)
T(sl, "低価格・高品質。用途に合わせたサイズバリエーションと\n寒冷地仕様を展開。",
  6.1, 1.95, 7.0, 0.75, sz=13, color=G1, wrap=True)
L(sl, 6.1, 2.85, 13.1, 2.85)
T(sl, "税抜価格", 6.1, 2.98, 7.0, 0.35, sz=12, bold=True, color=G2)
rows9 = [
    ("6m モデル",    "全長6000 × 全幅2380mm", "7,490,000"),
    ("7.2m モデル",  "全長7200 × 全幅2380mm", "7,990,000"),
    ("9m モデル",    "全長9000 × 全幅2380mm", "8,780,000"),
    ("寒冷地モデル", "高断熱仕様・片流れ屋根", "8,000,000"),
]
y9 = 3.38
for name, spec, price in rows9:
    R(sl, 6.1, y9, 6.95, 0.7, fc=RGBColor(0xF5,0xF5,0xF3))
    T(sl, name, 6.22, y9+0.08, 2.1, 0.3, sz=12, bold=True)
    T(sl, spec, 6.22, y9+0.38, 2.5, 0.25, sz=10, color=G2)
    T(sl, f"¥{price}〜", 8.7, y9+0.18, 4.2, 0.35, sz=15, bold=True, align=PP_ALIGN.RIGHT)
    y9 += 0.76
L(sl, 6.1, y9+0.1, 13.1, y9+0.1)
T(sl, "用途：グランピング宿泊棟 ／ サテライトオフィス ／ 店舗・受付 ／ セカンドハウス",
  6.1, y9+0.2, 7.0, 0.4, sz=11, color=G2)

# ═══════════════════════════════════════════════════════════════════
# 10 ROADIE & ROADIE mini
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "ROADIE & ROADIE mini", "自然と繋がるモデル")

PHOTO(sl, 0.45, 1.6, 6.0, 2.5, "ROADIE 外観写真")

T(sl, "ROADIE", 7.0, 1.65, 6.0, 0.65, sz=26, bold=True)
T(sl, "自然と一体化するシームレスデザイン。\nキャンプ場・リゾート施設に最適。",
  7.0, 2.35, 6.0, 0.7, sz=13, color=G1, wrap=True)
T(sl, "¥10,250,000〜", 7.0, 3.1, 5.5, 0.6, sz=24, bold=True)
T(sl, "寒冷地モデル　¥10,350,000〜", 7.0, 3.7, 5.5, 0.35, sz=13, color=G2)

L(sl, 0.45, 4.35, 12.85, 4.35, color=DIV)

PHOTO(sl, 0.45, 4.55, 2.9, 2.5, "ROADIE mini 平面図\n3.2m×2m / 6.51㎡")

T(sl, "ROADIE mini", 3.7, 4.6, 5.5, 0.6, sz=22, bold=True)
T(sl, "牽引免許不要の小型トレーラー（面積6.51㎡）。\n多目的にアレンジ可能な広々フラットタイプ。",
  3.7, 5.25, 5.5, 0.7, sz=13, color=G1, wrap=True)
T(sl, "¥3,700,000〜", 3.7, 6.0, 5.5, 0.6, sz=24, bold=True)

# ═══════════════════════════════════════════════════════════════════
# 11 MIGRA
# ═══════════════════════════════════════════════════════════════════
sl = S()
PHOTO(sl, 0.0, 0.0, 6.5, 7.5, "MIGRA 外観写真")
R(sl, 6.5, 0.0, 6.83, 7.5, fc=W)
R(sl, 6.5, 0.0, 0.06, 7.5, fc=INK)
T(sl, "MIGRA", 6.8, 0.45, 6.3, 0.85, sz=42, bold=True)
T(sl, "上質な木造シリーズ",  6.8, 1.35, 6.3, 0.45, sz=16, color=G2)
T(sl, "木の温かみを持つ快適なリビング空間。\n機能性と美しさを兼ね備えたハイエンドモデル。",
  6.8, 1.9, 6.3, 0.8, sz=13, color=G1, wrap=True)
L(sl, 6.8, 2.85, 13.1, 2.85)
rows_m = [
    ("標準モデル",              "¥9,950,000〜"),
    ("寒冷地モデル（W断熱/TG）", "¥10,050,000〜"),
    ("太陽光搭載（オフグリッド）","¥13,000,000〜"),
]
ym = 3.05
for name, price in rows_m:
    T(sl, name,  6.8, ym,       3.8, 0.45, sz=13, color=G1)
    T(sl, price, 9.8, ym-0.08, 3.1, 0.55, sz=20, bold=True, align=PP_ALIGN.RIGHT)
    ym += 0.72
L(sl, 6.8, ym+0.1, 13.1, ym+0.1)
T(sl, "用途：高級グランピング ／ プレミアムセカンドハウス ／ オフグリッド居住",
  6.8, ym+0.22, 6.3, 0.38, sz=11, color=G2)

# ═══════════════════════════════════════════════════════════════════
# 12 STRUCTURE & SPECS
# ═══════════════════════════════════════════════════════════════════
sl = S()
PHOTO(sl, 0.0, 0.0, 5.5, 7.5, "構造図 / 断面図")
R(sl, 5.5, 0.0, 7.83, 7.5, fc=W)
R(sl, 5.5, 0.0, 0.06, 7.5, fc=INK)

T(sl, "基本構造と性能", 5.8, 0.45, 7.2, 0.72, sz=28, bold=True)
T(sl, "長く安心して使える、プロ仕様の構造。", 5.8, 1.2, 7.2, 0.38, sz=13, color=G2)
L(sl, 5.8, 1.72, 13.1, 1.72)

specs = [
    ("堅牢な軽量鉄骨造",
     "移動・輸送の振動に耐える業務用設計。一般住宅と同等以上の強度。"),
    ("内装建材 F★★★★（フォースター）以上",
     "シックハウスの原因物質を最大限抑制。小さなお子様にも安心の室内環境。"),
    ("高性能断熱材",
     "床：ミラフォームMKS 40mm　天井：GW24K 100mm　壁：GW24K 50mm\n夏の猛暑・冬の寒冷地でも快適な室内環境を実現。"),
    ("耐久外装",
     "立平ロック25型 / 角波SK333採用。雨・風・紫外線に強くメンテコスト低減。"),
]
ys = 1.88
for ttl, body in specs:
    R(sl, 5.8, ys, 0.05, 0.85, fc=INK)   # left accent bar
    T(sl, ttl,  5.95, ys+0.04, 7.0, 0.38, sz=13, bold=True)
    T(sl, body, 5.95, ys+0.44, 7.0, 0.55, sz=11, color=G1, wrap=True)
    ys += 1.18

# ═══════════════════════════════════════════════════════════════════
# 13 CUSTOMIZATION
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "空間のカスタマイズ")

PHOTO(sl, 0.45, 1.6,  5.8, 2.25, "内観写真①")
PHOTO(sl, 0.45, 4.0,  5.8, 2.25, "内観写真②")

T(sl, "標準装備", 6.6, 1.6, 6.0, 0.38, sz=14, bold=True)
L(sl, 6.6, 2.05, 13.0, 2.05)
T(sl, "ダクトレール照明、ダウンライト、基本電気配線・スイッチ完備。",
  6.6, 2.15, 6.3, 0.5, sz=12, color=G1, wrap=True)

T(sl, "カスタムオプション", 6.6, 2.85, 6.0, 0.38, sz=14, bold=True)
L(sl, 6.6, 3.3, 13.0, 3.3)
opts = [
    ("水回り",   "トイレ、シャワールーム、洗面台、キッチン等"),
    ("空調設備", "エアコン設置（ブラケット含む）"),
    ("開口部",   "引き違いテラス窓、FIX窓、ドアの追加・位置変更"),
    ("内装仕上", "フロアタイル、塩ビタイル、シナ合板など"),
]
yo = 3.45
for k, v in opts:
    R(sl, 6.6, yo, 0.04, 0.38, fc=INK)
    T(sl, k, 6.75, yo+0.03, 1.4, 0.32, sz=12, bold=True)
    T(sl, v, 8.25, yo+0.03, 4.6, 0.32, sz=12, color=G1)
    yo += 0.5

# ═══════════════════════════════════════════════════════════════════
# 14 TRANSPORT PRICING
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "輸送料金について")
T(sl, "製造拠点から設置エリアまでの距離に応じた明瞭な輸送料金体系。",
  0.5, 1.55, 9.0, 0.38, sz=13)
T(sl, "（特殊条件・フェリー等は別途見積）", 0.5, 1.95, 6.0, 0.32, sz=11, color=G2)

series = [
    ("Tinys INSPIRATION", "製造地：名古屋"),
    ("ROADIE シリーズ",   "製造地：千葉"),
    ("MIGRA シリーズ",    "製造地：要確認"),
]
transport = [
    ("〜200km", "310,000"),
    ("〜500km", "570,000"),
    ("〜800km", "840,000"),
]
CW = 3.85
CX = [0.5, 4.5, 8.5]
for si, ((sname, mfg), cx) in enumerate(zip(series, CX)):
    R(sl, cx, 2.45, CW, 0.7, fc=INK)
    T(sl, sname, cx+0.12, 2.5,  CW-0.2, 0.38, sz=13, bold=True, color=W)
    T(sl, mfg,   cx+0.12, 2.88, CW-0.2, 0.25, sz=10, color=G3)
    for ri, (dist, price) in enumerate(transport):
        y = 3.25 + ri * 0.85
        fc = W if ri % 2 else RGBColor(0xF5,0xF5,0xF3)
        R(sl, cx, y, CW, 0.78, fc=fc)
        T(sl, dist, cx+0.12, y+0.22, 2.0, 0.35, sz=12, color=G1)
        T(sl, f"¥{price}", cx+2.0, y+0.15, CW-2.1, 0.48, sz=19, bold=True, align=PP_ALIGN.RIGHT)
T(sl, "※MIGRAシリーズは特殊車両許可が必要な場合があります。別途お見積もりをご案内します。",
  0.5, 5.8, 12.3, 0.32, sz=10, color=G3)

# ═══════════════════════════════════════════════════════════════════
# 15 REVENUE SIMULATION (new)
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "収益シミュレーション", "グランピング施設への1台導入モデルケース（ROADIE）")

# Left panel — cost
R(sl, 0.45, 1.6, 5.85, 5.55, fc=RGBColor(0xF5,0xF5,0xF3))
T(sl, "初期投資", 0.65, 1.75, 5.4, 0.4, sz=14, bold=True, color=G2)
invest = [
    ("ROADIE 本体価格",        "10,250,000"),
    ("輸送費（中距離）",        "570,000"),
    ("設置・造成工事（概算）",  "500,000"),
    ("カスタム・オプション",    "500,000"),
]
yi = 2.25
for name, amt in invest:
    T(sl, name,          0.65, yi, 3.6, 0.38, sz=12, color=G1)
    T(sl, f"¥{amt}", 3.6,  yi, 2.55, 0.38, sz=13, bold=True, align=PP_ALIGN.RIGHT)
    yi += 0.48
L(sl, 0.65, yi+0.05, 6.1, yi+0.05)
T(sl, "初期投資合計",    0.65, yi+0.15, 3.0, 0.5, sz=15, bold=True)
T(sl, "¥11,820,000", 3.0, yi+0.1,  2.9, 0.6, sz=22, bold=True, align=PP_ALIGN.RIGHT)

# Right panel — revenue & ROI
R(sl, 6.8, 1.6, 6.0, 5.55, fc=W)
R(sl, 6.8, 1.6, 6.0, 0.06, fc=INK)
T(sl, "年間収益シミュレーション", 7.0, 1.75, 5.6, 0.4, sz=14, bold=True, color=G2)
sims = [
    ("宿泊単価",  "¥25,000 ／泊"),
    ("稼 働 率",  "60%（年間 219泊）"),
]
ys2 = 2.28
for lbl, val in sims:
    T(sl, lbl, 7.0, ys2, 2.8, 0.4, sz=12, color=G1)
    T(sl, val, 9.1, ys2, 3.5, 0.4, sz=13, bold=True, align=PP_ALIGN.RIGHT)
    ys2 += 0.52
L(sl, 7.0, ys2+0.05, 12.6, ys2+0.05)
T(sl, "年間売上（概算）", 7.0, ys2+0.15, 3.5, 0.45, sz=14, bold=True)
T(sl, "¥5,475,000", 7.0, ys2+0.6, 5.55, 0.85, sz=38, bold=True, align=PP_ALIGN.RIGHT)
L(sl, 7.0, ys2+1.6, 12.6, ys2+1.6)
# ROI highlight
R(sl, 6.8, 5.25, 6.0, 1.6, fc=INK)
T(sl, "投資回収期間", 7.0, 5.38, 3.5, 0.38, sz=13, color=W)
T(sl, "約2年2ヶ月", 7.0, 5.75, 5.55, 0.85, sz=36, bold=True, color=W, align=PP_ALIGN.CENTER)
T(sl, "複数台導入でスケールメリット。中古売却でさらにコスト低減。",
  7.0, 6.62, 5.6, 0.32, sz=10, color=G3, align=PP_ALIGN.CENTER)
T(sl, "※上記はモデルケースです。実際の数値は立地・運営状況により異なります。",
  0.45, 7.22, 12.3, 0.25, sz=9, color=G3)

# ═══════════════════════════════════════════════════════════════════
# 16 SECOND-HAND PLATFORM
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "中古専門プラットフォーム")

T(sl, "TRAILER HOUSE SECOND HAND by YADOKARI",
  0.5, 1.58, 8.5, 0.6, sz=20, bold=True)
T(sl, "日本初の中古トレーラーハウス専門サイトを自社運営。\n導入後の売却（出口戦略）や、初期費用を抑えた中古購入を強力にサポート。",
  0.5, 2.25, 6.0, 0.85, sz=13, color=G1, wrap=True)
L(sl, 0.5, 3.25, 6.4, 3.25)
T(sl, "流通事例", 0.5, 3.4, 2.0, 0.35, sz=11, bold=True, color=G2)
examples = [
    ("コンパクトトレーラー（牽引免許不要）", "¥1,300,000"),
    ("宿泊向けトレーラーハウス（住箱）",     "¥3,500,000〜4,200,000"),
    ("オフィス・店舗向け 9m モデル",         "¥5,850,000"),
]
ye = 3.85
for name, price in examples:
    T(sl, "・" + name, 0.5,  ye, 4.0, 0.38, sz=12, color=G1)
    T(sl, price,       4.2,  ye, 2.1, 0.38, sz=13, bold=True, align=PP_ALIGN.RIGHT)
    ye += 0.48

# right: black showcase panel
R(sl, 6.8, 1.55, 6.0, 5.65, fc=INK)
T(sl, "TRAILER HOUSE\nSECOND HAND\nby YADOKARI",
  7.0, 2.6, 5.6, 2.2, sz=20, bold=True, color=W, align=PP_ALIGN.CENTER)
L(sl, 7.2, 4.9, 12.6, 4.9, color=G1, pt=0.5)
T(sl, "中古物件 掲載中", 7.0, 5.05, 5.6, 0.4, sz=13, color=G2, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# 17 PROCESS / FLOW
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "導入までの流れ")

steps = [
    ("ご相談・ヒアリング",    "1〜2 週間", "用途、予算、設置場所（法規制含む）の確認。"),
    ("プランニング・お見積り","2〜3 週間", "モデル選定・レイアウト作成・仕様決定。"),
    ("ご 契 約",              "1 週間",   "仕様確約および製造向け正式契約。"),
    ("製造・施工",            "8〜12 週間","国内工場で品質管理のもと製造・内装施工。"),
    ("納車・設置",            "1〜2 日",  "専門チームによる牽引納車・現地設置。"),
]
ys = 1.58
for i, (ttl, dur, body) in enumerate(steps):
    # Step number block
    R(sl, 0.45, ys, 0.6, 0.72, fc=INK)
    T(sl, str(i+1), 0.45, ys+0.1, 0.6, 0.52, sz=20, bold=True, color=W, align=PP_ALIGN.CENTER)
    # Title
    T(sl, ttl, 1.18, ys+0.1, 4.5, 0.45, sz=14, bold=True)
    # Duration badge
    R(sl, 5.8, ys+0.14, 1.9, 0.38, fc=RGBColor(0xEE,0xEE,0xEC))
    T(sl, dur, 5.82, ys+0.18, 1.86, 0.32, sz=11, color=G1, align=PP_ALIGN.CENTER)
    # Body
    T(sl, body, 7.9, ys+0.15, 5.1, 0.45, sz=12, color=G1, wrap=True)
    ys += 0.88
    if i < 4:
        L(sl, 0.45, ys-0.08, 13.0, ys-0.08, color=DIV)

# Total period bar
R(sl, 0.45, 6.7, 12.4, 0.65, fc=INK)
T(sl, "ご注文〜ご利用開始まで：約3〜4ヶ月",
  0.65, 6.78, 12.0, 0.5, sz=16, bold=True, color=W, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# 18 FAQ
# ═══════════════════════════════════════════════════════════════════
sl = S()
HBAR(sl, "よくあるご質問（FAQ）")

faqs = [
    ("Q1. 建築確認申請は必要ですか？",
     "一定の条件（随時かつ任意に移動できる状態等）を満たすことで「車両」扱いとなり、建築確認申請が不要になります。条件については個別にご確認ください。"),
    ("Q2. 固定資産税はかかりますか？",
     "車両として登録されている場合、土地への定着が認められないため原則として固定資産税の対象外となります。条件・自治体により異なる場合があります。"),
    ("Q3. ローン・リースは使えますか？",
     "各種金融機関のローンやリースに対応しています。詳細はご相談ください。"),
    ("Q4. 台風・大雪など災害時は大丈夫ですか？",
     "堅牢な軽量鉄骨造と高耐久外装材を採用し、一般建築物と同等以上の耐候性があります。大型台風時はアンカー固定を推奨します。"),
    ("Q5. メンテナンスサポートはありますか？",
     "定期点検・修繕対応・部品供給まで YADOKARIが一貫してサポートします。"),
    ("Q6. 使わなくなったらどうすればいいですか？",
     "「TRAILER HOUSE SECOND HAND by YADOKARI」にて中古売却が可能です。出口戦略までサポートするのが当社の強みです。"),
]
positions_faq = [
    (0.45, 1.6), (6.9, 1.6),
    (0.45, 3.45),(6.9, 3.45),
    (0.45, 5.3), (6.9, 5.3),
]
for (x, y), (q, a) in zip(positions_faq, faqs):
    R(sl, x, y, 0.04, 1.6, fc=INK)    # left accent bar
    T(sl, q, x+0.18, y+0.08, 5.9, 0.38, sz=12, bold=True)
    T(sl, a, x+0.18, y+0.52, 5.9, 0.98, sz=11, color=G1, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# 19 CLOSING
# ═══════════════════════════════════════════════════════════════════
sl = S(dark=True)
T(sl, "新しい暮らしと事業を、ここから。",
  1.0, 2.0, 11.3, 1.1, sz=38, bold=True, color=W, align=PP_ALIGN.CENTER)
L(sl, 3.8, 3.25, 9.5, 3.25, color=G1, pt=0.5)
T(sl, "トレーラーハウスに関するご相談、お見積りのご依頼はお気軽にご連絡ください。",
  1.0, 3.45, 11.3, 0.45, sz=14, color=G2, align=PP_ALIGN.CENTER)
T(sl, "WEB：https://yadokari.net/",
  1.0, 4.05, 11.3, 0.4, sz=13, color=G2, align=PP_ALIGN.CENTER)
T(sl, "お問い合わせ：YADOKARIオフィシャルサイト内「お問い合わせ」より",
  1.0, 4.55, 11.3, 0.4, sz=13, color=G2, align=PP_ALIGN.CENTER)
L(sl, 3.8, 5.1, 9.5, 5.1, color=G1, pt=0.5)
T(sl, "YADOKARI株式会社　｜　神奈川県横浜市保土ヶ谷区星川1-1-1 星天qlay 2階 qlaytion gallery",
  1.0, 7.1, 11.3, 0.32, sz=10, color=G2, align=PP_ALIGN.CENTER)

# ── save ────────────────────────────────────────────────────
out = "/home/user/test-labs/YADOKARI_Proposal_v3_Branded.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
