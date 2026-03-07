from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x99, 0x99, 0x99)
LGRAY      = RGBColor(0xDD, 0xDD, 0xDD)
DGRAY      = RGBColor(0x33, 0x33, 0x33)
ACCENT     = RGBColor(0x1A, 0x1A, 0x1A)

# ── helpers ──────────────────────────────────────────────────────────────────

def S():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=WHITE):
    from pptx.util import Inches
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color=LGRAY, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape

def line(slide, x1, y1, x2, y2, color=LGRAY, width_pt=0.75):
    from pptx.util import Pt as UPt
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = UPt(width_pt)
    return c

def tb(slide, text, l, t, w, h,
       size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def mtb(slide, l, t, w, h, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    return tf

def ap(tf, text, size=13, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, space_before_pt=0, space_after_pt=0):
    """Add paragraph to existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_before_pt:
        p.space_before = Pt(space_before_pt)
    if space_after_pt:
        p.space_after = Pt(space_after_pt)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def header(slide, title, subtitle=None):
    tb(slide, title, 0.5, 0.25, 8, 0.75, size=28, bold=True)
    if subtitle:
        tb(slide, subtitle, 0.5, 1.0, 8, 0.4, size=14, color=GRAY)
    line(slide, 0.5, 1.05 if subtitle else 1.05, 12.8, 1.05 if subtitle else 1.05)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 ── タイトル
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
tb(sl, "小さいほど、豊かだ。", 0.5, 0.6, 6, 1.2, size=40, bold=True)
line(sl, 0.5, 1.9, 5.5, 1.9)
tb(sl, "YADOKARI Trailer House Proposal", 0.5, 2.1, 6, 0.5, size=18, color=DGRAY)
tb(sl, "YADOKARI株式会社", 0.5, 6.8, 4, 0.4, size=12, color=GRAY)
rect(sl, 6.5, 0.3, 6.5, 6.9, fill_color=RGBColor(0xEE,0xEE,0xEE))
tb(sl, "［表紙写真エリア］\nトレーラーハウス外観写真を配置してください",
   6.7, 3.2, 6.1, 0.8, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 ── ミッション + 実績
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
line(sl, 6.5, 0.5, 6.5, 7.0, color=LGRAY)

tf = mtb(sl, 0.5, 1.5, 5.5, 4.0)
ap(tf, "世界を変える、", size=34, bold=True)
ap(tf, "暮らしを創る。", size=34, bold=True)

tf2 = mtb(sl, 7.0, 1.2, 5.8, 5.5)
ap(tf2, "「世界を変える、暮らしを創る」の実現を", size=14, space_before_pt=0)
ap(tf2, "目指すソーシャルデザインカンパニー。", size=14, space_after_pt=6)
ap(tf2, "小さな空間の中に広がる、大きな可能性。", size=14, space_after_pt=6)
ap(tf2, "多様な生き方と暮らしのカタチに触れ、", size=14)
ap(tf2, "未来のコライフスタイルを探求し、", size=14)
ap(tf2, "提供し続けます。", size=14, space_after_pt=16)
line(sl, 7.0, 4.2, 12.8, 4.2, color=LGRAY)
ap(tf2, "実績", size=13, bold=True, color=GRAY, space_before_pt=8)
ap(tf2, "販売実績　　●●台以上", size=13, space_before_pt=4)
ap(tf2, "導入施設　　●●件以上", size=13)
ap(tf2, "メディア掲載　●●誌 / NHK 他●媒体", size=13)
ap(tf2, "設立　　　　2012年 / 神奈川県横浜市", size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 ── 【新規】市場背景・なぜ今か
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "なぜ今、トレーラーハウスなのか？", "市場背景")

cols = [
    ("インバウンド回復", "訪日観光客の急増に伴い、\nグランピング・体験型宿泊の\n需要が急拡大。"),
    ("移住・二拠点生活", "地方移住や二拠点生活トレンドが\n加速。手軽に始められる\n住まいのニーズが増大。"),
    ("遊休地の有効活用", "空き地・農地・山林などの\n遊休資産を収益化したい\nオーナーが増加。"),
    ("建設業の課題", "職人不足・工期長期化・\nコスト高騰が深刻化。\nスピードと低コストが求められる。"),
]
for i, (ttl, body) in enumerate(cols):
    x = 0.5 + i * 3.2
    rect(sl, x, 1.4, 2.9, 0.55, fill_color=DGRAY)
    tb(sl, ttl, x+0.1, 1.45, 2.7, 0.45, size=13, bold=True, color=WHITE)
    tb(sl, body, x+0.1, 2.05, 2.7, 1.8, size=12, color=DGRAY, wrap=True)

line(sl, 0.5, 4.1, 12.8, 4.1)
tf = mtb(sl, 0.5, 4.3, 12.3, 2.8)
ap(tf, "YADOKARIのトレーラーハウスは、これらすべての課題に応える「可動産」という新しい選択肢です。",
   size=16, bold=True, space_after_pt=10)
ap(tf, "・初期費用を抑えながら高品質な空間を短期間で実現　　・移動できるから、土地を選ばない柔軟な事業展開が可能",
   size=13, color=DGRAY)
ap(tf, "・建築物でない「車両」扱いで、法規制・税制面のハードルを低減　　・中古売却（出口戦略）まで一貫サポート",
   size=13, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 ── 可動産という選択肢（改訂：比較表追加）
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "可動産という選択肢")

icons_data = [
    ("1. 機動力", "基礎を必要とせず、\n場所にとらわれない\n柔軟な運用が可能。"),
    ("2. コストパフォーマンス", "従来の建築物と比較して\n初期費用を大幅に抑制。"),
    ("3. スピード導入", "製造から納品までの\nリードタイムが短く、\n素早い事業展開が可能。"),
    ("4. 法規制のメリット", "一定の条件を満たすことで\n「車両」扱いとなり、\n税制・設置のハードルを低減。"),
]
for i, (ttl, body) in enumerate(icons_data):
    x = 0.5 + i * 3.2
    if i > 0:
        line(sl, x - 0.1, 1.2, x - 0.1, 3.8, color=LGRAY)
    tb(sl, ttl, x, 1.3, 2.9, 0.5, size=14, bold=True)
    tb(sl, body, x, 1.9, 2.9, 1.8, size=12, color=DGRAY, wrap=True)

line(sl, 0.5, 3.9, 12.8, 3.9)
tb(sl, "従来の建築物との比較", 0.5, 4.1, 4, 0.4, size=13, bold=True, color=GRAY)

headers_row = ["", "従来の建築物", "トレーラーハウス"]
rows = [
    ["初期費用", "数千万〜数億円", "749万円〜"],
    ["工  期", "6ヶ月〜1年以上", "約3〜4ヶ月"],
    ["移設・撤去", "ほぼ不可能", "可能（再設置・売却）"],
    ["建築確認申請", "必要", "不要（条件付き）※"],
    ["固定資産税", "課税", "非課税（条件付き）※"],
    ["出口戦略", "解体コスト発生", "中古売却が可能"],
]
col_w = [2.8, 3.5, 3.5]
col_x = [0.5, 3.4, 7.0]
row_h = 0.33
row_y_start = 4.55

for ci, hd in enumerate(headers_row):
    rect(sl, col_x[ci], row_y_start, col_w[ci] - 0.05, row_h,
         fill_color=DGRAY if ci > 0 else RGBColor(0xF5,0xF5,0xF5))
    tb(sl, hd, col_x[ci]+0.05, row_y_start+0.04, col_w[ci]-0.1, row_h-0.05,
       size=11, bold=True, color=WHITE if ci > 0 else GRAY, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    y = row_y_start + row_h + ri * row_h
    for ci, cell in enumerate(row):
        fill = RGBColor(0xF8,0xF8,0xF8) if ri % 2 == 0 else WHITE
        if ci == 2:
            fill = RGBColor(0xF0,0xF5,0xF0)
        rect(sl, col_x[ci], y, col_w[ci] - 0.05, row_h - 0.02, fill_color=fill)
        tb(sl, cell, col_x[ci]+0.08, y+0.04, col_w[ci]-0.15, row_h-0.06,
           size=11, bold=(ci==2), align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

tb(sl, "※一定条件下において。詳細はご相談ください。", 0.5, 7.15, 8, 0.3, size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 ── 多様な活用シーン
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "多様な活用シーン")

scenes = [
    ("宿泊施設・グランピング", "高いデザイン性で宿泊単価を向上。\n複数台導入によるヴィレッジ展開にも最適。"),
    ("店舗・カフェ", "洗練された外観で集客力を強化。\nスモールスタートでの新規事業に。"),
    ("オフィス・事務所", "節税対策や、自然の中での\nサテライトオフィス構築に。"),
    ("住まい・セカンドハウス", "趣味の空間や、\n二拠点居住のベースキャンプとして。"),
]
positions = [(0.4, 1.2), (6.9, 1.2), (0.4, 4.2), (6.9, 4.2)]
for i, ((x, y), (ttl, body)) in enumerate(zip(positions, scenes)):
    rect(sl, x, y, 6.1, 1.8, fill_color=RGBColor(0xEE,0xEE,0xEE))
    tb(sl, "［写真エリア］", x+2.3, y+0.75, 2.0, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    tb(sl, ttl, x, y+1.85, 6.1, 0.4, size=14, bold=True)
    tb(sl, body, x, y+2.3, 6.1, 0.8, size=12, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 ── YADOKARIの強み
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "YADOKARIの強み")

strengths = [
    ("洗練された\nデザイン",
     "空間の豊かさを追求した高いデザイン性。\n用途に合わせたセミオーダー提案が可能。"),
    ("圧倒的な\nコストパフォーマンス",
     "無駄を省いた設計と標準化による、\n初期投資の最小化と高いROIの実現。"),
    ("一気通貫の\nエコシステム",
     "企画・設計から納品、さらには業界初の\n中古売買プラットフォーム（二次流通）\nまで完全サポート。"),
]
for i, (ttl, body) in enumerate(strengths):
    x = 0.5 + i * 4.3
    if i > 0:
        line(sl, x - 0.15, 1.2, x - 0.15, 6.5, color=LGRAY)
    tb(sl, ttl, x, 2.0, 4.0, 0.9, size=20, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, body, x, 3.2, 4.0, 2.5, size=13, color=DGRAY, align=PP_ALIGN.CENTER, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 ── 【新規】導入事例
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "導入事例", "実際の活用シーンと成果")

cases = [
    {
        "title": "事例① グランピング施設（地方リゾート）",
        "model": "ROADIE × 5台",
        "issue": "課題：短期間で宿泊棟を増やしたかった",
        "results": [
            "工期4ヶ月で5棟完成（従来比1/3の期間）",
            "導入後6ヶ月で稼働率80%達成",
            "宿泊単価を従来の1.5倍に設定",
        ],
        "voice": "「建築物では実現できなかったスピードと\nデザイン性を両立できました。」\n── ○○施設 オーナー様",
    },
    {
        "title": "事例② サテライトオフィス（首都圏近郊）",
        "model": "Tinys INSPIRATION 9m × 2台",
        "issue": "課題：テレワーク対応の執務スペースが必要",
        "results": [
            "敷地内に2ヶ月で設置完了",
            "固定資産税の節税効果あり",
            "社員満足度が向上、採用PRにも活用",
        ],
        "voice": "「こんなに早く、コストを抑えて\n快適な空間が手に入るとは思いませんでした。」\n── ○○株式会社 担当者様",
    },
]

for i, c in enumerate(cases):
    x = 0.5 + i * 6.4
    rect(sl, x, 1.3, 6.1, 5.9, fill_color=RGBColor(0xF8,0xF8,0xF8))
    rect(sl, x, 1.3, 6.1, 1.5, fill_color=DGRAY)
    tb(sl, c["title"], x+0.15, 1.35, 5.8, 0.45, size=13, bold=True, color=WHITE)
    tb(sl, c["model"], x+0.15, 1.8, 5.8, 0.35, size=12, color=LGRAY)
    tb(sl, c["issue"], x+0.15, 2.9, 5.8, 0.35, size=12, bold=True, color=DGRAY)
    tf = mtb(sl, x+0.15, 3.35, 5.8, 1.5)
    for r in c["results"]:
        ap(tf, "✓ " + r, size=12, color=DGRAY, space_before_pt=2)
    line(sl, x+0.1, 5.05, x+5.9, 5.05, color=LGRAY)
    tb(sl, c["voice"], x+0.15, 5.15, 5.8, 1.9, size=11, color=DGRAY, wrap=True)

tb(sl, "※事例写真・詳細データは別途ご用意できます。実績数値はすべてモデルケースです。",
   0.5, 7.2, 12.3, 0.3, size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 ── プロダクトラインナップ
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "プロダクトラインナップ")

models = [
    ("Tinys INSPIRATION",
     "低価格・高品質・様々なバリエーションを\n取り揃えたベーシックシリーズ。",
     "¥7,490,000〜"),
    ("ROADIE",
     "周囲の自然とシームレスに繋がるデザインで、\n自然の中の暮らしを実現するモデル。",
     "¥10,250,000〜"),
    ("MIGRA",
     "木質感と開放的なリビングの中で\n上質な時間を過ごすためのプレミアム木造シリーズ。",
     "¥9,950,000〜"),
]
for i, (name, desc, price) in enumerate(models):
    x = 0.5 + i * 4.3
    if i > 0:
        line(sl, x - 0.15, 1.2, x - 0.15, 7.0, color=LGRAY)
    rect(sl, x, 1.3, 4.0, 2.5, fill_color=RGBColor(0xE8,0xE8,0xE8))
    tb(sl, "［製品写真エリア］", x+1.0, 2.35, 2.0, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    tb(sl, name, x, 3.95, 4.0, 0.5, size=16, bold=True)
    tb(sl, desc, x, 4.55, 4.0, 1.2, size=12, color=DGRAY, wrap=True)
    tb(sl, price, x, 5.9, 4.0, 0.5, size=15, bold=True, color=ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 ── Tinys INSPIRATION 詳細
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
rect(sl, 0.0, 0.0, 6.2, 7.5, fill_color=RGBColor(0xEE,0xEE,0xEE))
tb(sl, "［製品写真エリア］", 1.5, 3.4, 3.0, 0.5, size=11, color=GRAY, align=PP_ALIGN.CENTER)

tb(sl, "Tinys INSPIRATION", 6.5, 0.5, 6.5, 0.7, size=30, bold=True)
tb(sl, "定番ベーシックモデル", 6.5, 1.25, 6.5, 0.4, size=14, color=GRAY)
line(sl, 6.5, 1.75, 12.8, 1.75)
tb(sl, "低価格・高品質。用途に合わせたサイズバリエーションと寒冷地仕様も展開。",
   6.5, 1.9, 6.3, 0.7, size=15, bold=True, wrap=True)
line(sl, 6.5, 2.75, 12.8, 2.75)
tb(sl, "税抜価格", 6.5, 2.9, 6.3, 0.4, size=14, bold=True)

tf = mtb(sl, 6.5, 3.35, 6.3, 3.5)
rows9 = [
    ("6mモデル", "全長6000×全幅2380mm", "¥7,490,000〜"),
    ("7.2mモデル", "全長7200×全幅2380mm", "¥7,990,000〜\n（※限定即納車：¥8,150,000）"),
    ("9mモデル", "全長9000×全幅2380mm", "¥8,780,000〜\n（※限定即納車：¥9,500,000）"),
    ("寒冷地モデル", "高断熱仕様・片流れ屋根", "¥8,000,000〜"),
]
for name, spec, price in rows9:
    ap(tf, f"・{name}：{spec}　｜　{price}", size=12, space_before_pt=5)

line(sl, 6.5, 6.3, 12.8, 6.3)
tb(sl, "■ こんな用途に最適", 6.5, 6.4, 6.3, 0.3, size=12, bold=True, color=GRAY)
tb(sl, "グランピング宿泊棟 ／ サテライトオフィス ／ 店舗・受付 ／ セカンドハウス",
   6.5, 6.75, 6.3, 0.45, size=11, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 ── ROADIE & ROADIE mini
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "ROADIE & ROADIE mini", "自然と繋がるモデル")

rect(sl, 0.4, 1.3, 6.1, 2.8, fill_color=RGBColor(0xE8,0xE8,0xE8))
tb(sl, "［ROADIE 写真エリア］", 1.8, 2.55, 2.8, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)

line(sl, 6.8, 1.2, 6.8, 7.2, color=LGRAY)

tb(sl, "ROADIE", 7.0, 1.3, 5.8, 0.5, size=18, bold=True)
tf = mtb(sl, 7.0, 1.85, 5.8, 2.0)
ap(tf, "・キャンプ場やリゾート施設に最適。豊かな自然を", size=12)
ap(tf, "　借景にするシームレスデザイン。", size=12)
ap(tf, "・標準モデル：¥10,250,000〜", size=12, space_before_pt=5)
ap(tf, "・寒冷地モデル：¥10,350,000〜", size=12)

line(sl, 0.4, 4.25, 12.8, 4.25, color=LGRAY)

rect(sl, 0.4, 4.45, 2.8, 2.5, fill_color=RGBColor(0xE8,0xE8,0xE8))
tb(sl, "平面図\n（約3.2m × 約2m\n 面積6.51㎡）",
   0.5, 5.35, 2.5, 1.0, size=10, color=GRAY, align=PP_ALIGN.CENTER)

tb(sl, "ROADIE mini", 7.0, 4.5, 5.8, 0.5, size=18, bold=True)
tf2 = mtb(sl, 7.0, 5.05, 5.8, 2.0)
ap(tf2, "・牽引免許不要の小型トレーラー（面積6.51㎡）。", size=12)
ap(tf2, "　多目的にアレンジ可能な広々フラットタイプ。", size=12)
ap(tf2, "・標準モデル：¥3,700,000〜", size=12, space_before_pt=5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 ── MIGRA
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
rect(sl, 0.0, 0.0, 6.5, 7.5, fill_color=RGBColor(0xE8,0xE8,0xE8))
tb(sl, "［MIGRA 写真エリア］", 1.8, 3.4, 2.8, 0.5, size=11, color=GRAY, align=PP_ALIGN.CENTER)

tb(sl, "MIGRA", 6.8, 0.5, 6.2, 0.7, size=36, bold=True)
tb(sl, "上質な木造シリーズ", 6.8, 1.25, 6.2, 0.45, size=16, bold=True, color=DGRAY)
tb(sl, "木の温かみを持つ快適なリビング空間。\n機能性と美しさを兼ね備えたハイエンドトレーラーハウス。",
   6.8, 1.85, 6.2, 0.9, size=14, color=DGRAY, wrap=True)
line(sl, 6.8, 2.95, 12.8, 2.95)

tf = mtb(sl, 6.8, 3.15, 6.2, 3.0)
ap(tf, "・標準モデル：¥9,950,000〜", size=13, space_before_pt=6)
ap(tf, "・寒冷地モデル（W断熱・トリプルガラス）：¥10,050,000〜", size=13, space_before_pt=6)
ap(tf, "・太陽光搭載モデル（オフグリッド対応）：¥13,000,000〜", size=13, space_before_pt=6)

line(sl, 6.8, 5.3, 12.8, 5.3)
tb(sl, "■ こんな用途に最適", 6.8, 5.45, 6.2, 0.3, size=12, bold=True, color=GRAY)
tb(sl, "高級グランピング ／ プレミアムセカンドハウス ／ オフグリッド居住",
   6.8, 5.8, 6.2, 0.4, size=11, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 ── 基本構造と性能（顧客目線改訂）
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
rect(sl, 0.0, 0.0, 5.8, 7.5, fill_color=RGBColor(0xF0,0xF0,0xF0))
tb(sl, "［構造図エリア］", 1.5, 3.4, 2.8, 0.5, size=11, color=GRAY, align=PP_ALIGN.CENTER)

tb(sl, "基本構造と性能", 6.2, 0.5, 6.8, 0.65, size=26, bold=True)
tb(sl, "長く安心して使える、プロ仕様の構造。", 6.2, 1.2, 6.8, 0.4, size=14, color=GRAY)
line(sl, 6.2, 1.7, 12.8, 1.7)

specs = [
    ("構造：堅牢な軽量鉄骨造",
     "移動・輸送の振動に耐える業務用設計。一般住宅と同等以上の強度。"),
    ("安全性：内装建材 F★★★★（フォースター）以上",
     "シックハウス症候群の原因物質を最大限に抑制。\n小さなお子様にも安心の室内環境。"),
    ("断熱性能：床・天井・壁に高性能断熱材を採用",
     "床：ミラフォームMKS 40mm　天井：GW24K 100mm　壁：GW24K 50mm\n夏の猛暑・冬の寒冷地でも快適な室内環境を実現。"),
    ("外装：耐久性に優れた立平ロック25型 / 角波SK333",
     "雨・風・紫外線に強く、メンテナンスコストを最小化。"),
]
y = 1.9
for ttl, body in specs:
    tb(sl, "・" + ttl, 6.2, y, 6.6, 0.4, size=13, bold=True)
    tb(sl, body, 6.4, y + 0.42, 6.4, 0.65, size=11, color=DGRAY, wrap=True)
    y += 1.2
    line(sl, 6.2, y - 0.08, 12.8, y - 0.08, color=LGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 ── 空間のカスタマイズ
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "空間のカスタマイズ")

rect(sl, 0.4, 1.3, 5.5, 2.3, fill_color=RGBColor(0xE8,0xE8,0xE8))
tb(sl, "［内観写真①］", 2.0, 2.3, 2.3, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)
rect(sl, 0.4, 3.75, 5.5, 2.3, fill_color=RGBColor(0xE8,0xE8,0xE8))
tb(sl, "［内観写真②］", 2.0, 4.75, 2.3, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)

tb(sl, "標準装備：", 6.3, 1.3, 6.5, 0.4, size=14, bold=True)
line(sl, 6.3, 1.75, 12.8, 1.75)
tb(sl, "ダクトレール照明、ダウンライト、基本電気配線・スイッチ完備。",
   6.3, 1.85, 6.5, 0.5, size=12, color=DGRAY)

tb(sl, "カスタムオプション：", 6.3, 2.55, 6.5, 0.4, size=14, bold=True)
line(sl, 6.3, 3.0, 12.8, 3.0)

opts = [
    ("水回り", "トイレ、シャワールーム、洗面台、キッチン等のインフラ実装"),
    ("空調設備", "エアコン設置（ブラケット含む）"),
    ("開口部", "引き違いテラス窓、FIX窓、各種ドアの追加・位置変更"),
    ("内装仕上", "フロアタイル、塩ビタイル、シナ合板など用途に合わせたマテリアル変更"),
]
tf = mtb(sl, 6.3, 3.1, 6.5, 3.5)
for k, v in opts:
    ap(tf, f"・{k}：{v}", size=12, space_before_pt=5, color=DGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 ── 輸送料金（MIGRA追加）
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "輸送料金について")

tb(sl, "製造拠点から設置エリアまでの距離に応じた明瞭な輸送料金体系。",
   0.5, 1.2, 8.5, 0.4, size=13)
tb(sl, "（※特殊条件・フェリー等は要別途見積）", 0.5, 1.65, 8.5, 0.35, size=12, color=GRAY)

series_data = [
    ("Tinys INSPIRATIONシリーズ", "製造地：名古屋"),
    ("ROADIEシリーズ", "製造地：千葉"),
    ("MIGRAシリーズ", "製造地：要確認"),
]
transport_rows = [
    ("短距離（〜200km）", "¥310,000"),
    ("中距離（〜500km）", "¥570,000"),
    ("遠距離（〜800km）", "¥840,000"),
]
col_start = [0.5, 4.55, 8.6]
for si, ((series, mfg), cx) in enumerate(zip(series_data, col_start)):
    rect(sl, cx, 2.15, 3.85, 0.55, fill_color=DGRAY)
    tf = mtb(sl, cx + 0.1, 2.2, 3.65, 0.45)
    ap(tf, series, size=12, bold=True, color=WHITE)
    ap(tf, mfg, size=10, color=LGRAY)
    line(sl, cx, 2.7, cx + 3.85, 2.7, color=LGRAY)
    for ri, (dist, price) in enumerate(transport_rows):
        y = 2.82 + ri * 0.7
        fill = RGBColor(0xF8,0xF8,0xF8) if ri % 2 == 0 else WHITE
        rect(sl, cx, y, 3.85, 0.65, fill_color=fill)
        tb(sl, dist, cx + 0.1, y + 0.15, 2.2, 0.35, size=11)
        tb(sl, price, cx + 2.4, y + 0.15, 1.3, 0.35, size=12, bold=True, align=PP_ALIGN.RIGHT)

tb(sl, "※MIGRAシリーズは大型モデルのため、特殊車両許可が必要な場合があります。別途お見積もりをご案内します。",
   0.5, 5.4, 12.3, 0.4, size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 ── 【新規】収益シミュレーション
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "収益シミュレーション", "グランピング施設への1台導入モデルケース（ROADIE）")

# 左：投資
rect(sl, 0.5, 1.3, 5.8, 5.5, fill_color=RGBColor(0xF5,0xF5,0xF5))
tb(sl, "【初期投資】", 0.7, 1.5, 5.4, 0.4, size=14, bold=True)
line(sl, 0.7, 1.95, 6.1, 1.95, color=LGRAY)
tf = mtb(sl, 0.7, 2.1, 5.4, 3.5)
items = [
    ("ROADIE 本体価格", "¥10,250,000"),
    ("輸送費（中距離想定）", "¥570,000"),
    ("設置・造成工事（概算）", "¥500,000"),
    ("カスタム・オプション（概算）", "¥500,000"),
]
for name, amt in items:
    p = ap(tf, name, size=12, space_before_pt=4)
    # price on same line via spaces trick
    tf2_inner = mtb(sl, 0.7, 0, 5.4, 0.1)  # dummy, we'll use the text approach
for name, amt in items:
    pass
# Use a table-like layout
y_inv = 2.1
for name, amt in items:
    tb(sl, name, 0.8, y_inv, 3.8, 0.38, size=12, color=DGRAY)
    tb(sl, amt, 4.2, y_inv, 1.7, 0.38, size=12, align=PP_ALIGN.RIGHT)
    y_inv += 0.44
line(sl, 0.7, y_inv + 0.05, 6.1, y_inv + 0.05, color=BLACK)
tb(sl, "初期投資合計", 0.8, y_inv + 0.12, 3.8, 0.45, size=14, bold=True)
tb(sl, "¥11,820,000", 3.5, y_inv + 0.12, 2.4, 0.45, size=14, bold=True, align=PP_ALIGN.RIGHT)

# 右：収益
rect(sl, 6.8, 1.3, 6.0, 5.5, fill_color=RGBColor(0xF0,0xF5,0xF0))
tb(sl, "【年間収益シミュレーション】", 7.0, 1.5, 5.6, 0.4, size=14, bold=True)
line(sl, 7.0, 1.95, 12.6, 1.95, color=LGRAY)
sim_items = [
    ("宿泊単価", "¥25,000 ／泊"),
    ("稼働率", "60%（年間219泊）"),
    ("年間売上", "約¥5,475,000"),
]
y_sim = 2.15
for name, val in sim_items:
    tb(sl, name, 7.1, y_sim, 3.0, 0.4, size=12, color=DGRAY)
    tb(sl, val, 9.5, y_sim, 3.0, 0.4, size=12, bold=(name=="年間売上"), align=PP_ALIGN.RIGHT)
    y_sim += 0.5

line(sl, 7.0, y_sim + 0.1, 12.6, y_sim + 0.1, color=BLACK)
rect(sl, 6.8, y_sim + 0.2, 6.0, 1.1, fill_color=DGRAY)
tb(sl, "投資回収期間", 7.0, y_sim + 0.3, 3.0, 0.4, size=14, bold=True, color=WHITE)
tb(sl, "約2年2ヶ月", 9.8, y_sim + 0.3, 2.8, 0.4, size=20, bold=True, color=RGBColor(0x90,0xEE,0x90), align=PP_ALIGN.RIGHT)
tb(sl, "※複数台導入でスケールメリットあり。中古売却（出口）により実質コストをさらに低減可能。",
   7.0, y_sim + 1.45, 5.7, 0.45, size=10, color=DGRAY, wrap=True)

tb(sl, "※上記はモデルケースです。実際の数値は立地・運営状況により異なります。詳細シミュレーションはご相談ください。",
   0.5, 7.0, 12.3, 0.35, size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 ── 中古専門プラットフォーム
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "中古専門プラットフォーム")

tf = mtb(sl, 0.5, 1.3, 5.8, 2.5)
ap(tf, "「TRAILER HOUSE SECOND HAND by YADOKARI」。", size=14, bold=True)
ap(tf, "日本初の中古トレーラーハウス専門サイトを自社運営。", size=13, space_before_pt=6)
ap(tf, "導入後の売却（出口戦略）や、初期費用を抑えた", size=13, space_before_pt=4)
ap(tf, "中古購入を強力にサポート。", size=13)

line(sl, 0.5, 3.9, 6.1, 3.9, color=LGRAY)
tb(sl, "流通事例：", 0.5, 4.05, 5.5, 0.4, size=13, bold=True)
examples = [
    "牽引免許不要コンパクトトレーラー：¥1,300,000",
    "宿泊向けトレーラーハウス（住箱）：¥3,500,000〜¥4,200,000",
    "オフィス・店舗向け9mモデル：¥5,850,000",
]
tf2 = mtb(sl, 0.5, 4.5, 5.8, 2.0)
for ex in examples:
    ap(tf2, "・" + ex, size=12, space_before_pt=5, color=DGRAY)

rect(sl, 6.8, 1.3, 6.0, 5.8, fill_color=RGBColor(0x22,0x22,0x22))
tb(sl, "TRAILER HOUSE\nSECOND HAND\nby YADOKARI",
   7.0, 2.8, 5.6, 1.8, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "中古物件掲載中", 7.0, 4.8, 5.6, 0.5, size=13, color=LGRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 ── 導入までの流れ（期間追加）
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "導入までの流れ")

steps = [
    ("1. ご相談・ヒアリング", "約1〜2週間",
     "用途、予算、設置場所の条件（法規制含む）の確認。"),
    ("2. プランニング・お見積り", "約2〜3週間",
     "モデル選定、レイアウト作成、カスタム仕様の決定と精緻なお見積り。"),
    ("3. ご契約", "約1週間",
     "仕様確約および製造に向けた正式契約。"),
    ("4. 製造・施工", "約8〜12週間",
     "徹底した品質管理のもと、国内工場にて製造・内装施工。"),
    ("5. 納車・設置", "1〜2日",
     "専門チームによる牽引納車、指定場所への確実な設置作業。"),
]
line(sl, 6.5, 1.0, 6.5, 7.2, color=LGRAY)
y = 1.3
for i, (ttl, duration, body) in enumerate(steps):
    # Left: step title + duration
    rect(sl, 0.5, y, 1.1, 0.55, fill_color=DGRAY)
    tb(sl, str(i+1), 0.5, y + 0.05, 1.1, 0.45, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, ttl, 1.75, y + 0.04, 4.4, 0.4, size=14, bold=True)
    rect(sl, 4.6, y + 0.05, 1.6, 0.38, fill_color=RGBColor(0xEE,0xEE,0xEE))
    tb(sl, duration, 4.65, y + 0.1, 1.5, 0.3, size=11, color=DGRAY, align=PP_ALIGN.CENTER)
    # Right: description
    tb(sl, body, 6.8, y + 0.05, 6.0, 0.7, size=12, color=DGRAY, wrap=True)
    y += 1.0
    if i < 4:
        line(sl, 0.5, y - 0.1, 12.8, y - 0.1, color=LGRAY)

rect(sl, 0.5, 6.65, 12.3, 0.6, fill_color=DGRAY)
tb(sl, "ご注文〜ご利用開始まで：約3〜4ヶ月",
   0.7, 6.72, 12.0, 0.45, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 ── 【新規】FAQ
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl)
header(sl, "よくあるご質問（FAQ）")

faqs = [
    ("Q1. 建築確認申請は必要ですか？",
     "一定の条件（随時かつ任意に移動できる状態等）を満たすことで「車両」扱いとなり、建築確認申請が不要になります。設置条件については個別にご確認ください。"),
    ("Q2. 固定資産税はかかりますか？",
     "車両として登録されている場合、土地への定着が認められないため原則として固定資産税の対象外となります。※条件・自治体により異なる場合があります。"),
    ("Q3. ローン・リースは使えますか？",
     "各種金融機関のローンやリースに対応しています。詳細はご相談ください。"),
    ("Q4. 台風・大雪など災害時は大丈夫ですか？",
     "堅牢な軽量鉄骨造と高耐久外装材を採用。一般建築物と同等以上の耐候性があります。※大型台風時はアンカー固定を推奨します。"),
    ("Q5. メンテナンスサポートはありますか？",
     "定期点検・修繕対応・部品供給まで YADOKARIが一貫してサポートします。"),
    ("Q6. 使わなくなったらどうすればいいですか？",
     "「TRAILER HOUSE SECOND HAND by YADOKARI」にて中古売却が可能です。出口戦略までサポートするのが当社の強みです。"),
]
positions_faq = [
    (0.5, 1.25), (6.9, 1.25),
    (0.5, 3.1),  (6.9, 3.1),
    (0.5, 4.95), (6.9, 4.95),
]
for (x, y), (q, a) in zip(positions_faq, faqs):
    rect(sl, x, y, 6.1, 1.65, fill_color=RGBColor(0xF5,0xF5,0xF5))
    tb(sl, q, x + 0.15, y + 0.1, 5.8, 0.45, size=12, bold=True, color=DGRAY)
    tb(sl, a, x + 0.15, y + 0.58, 5.8, 0.98, size=11, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 ── クロージング
# ═══════════════════════════════════════════════════════════════════════════════
sl = S(); bg(sl, color=RGBColor(0xF8,0xF8,0xF8))
tb(sl, "新しい暮らしと事業を、ここから。",
   1.0, 2.5, 11.3, 1.0, size=34, bold=True, align=PP_ALIGN.CENTER)
line(sl, 3.5, 3.65, 9.8, 3.65)
tb(sl, "トレーラーハウスに関するご相談、お見積りのご依頼はお気軽にご連絡ください。",
   1.0, 3.85, 11.3, 0.5, size=14, color=DGRAY, align=PP_ALIGN.CENTER)
tb(sl, "WEB: https://yadokari.net/",
   1.0, 4.45, 11.3, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)
tb(sl, "お問い合わせフォーム：YADOKARIオフィシャルサイト内「お問い合わせ」より",
   1.0, 4.95, 11.3, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)
line(sl, 3.5, 5.55, 9.8, 5.55)
tb(sl, "YADOKARI株式会社　｜　神奈川県横浜市保土ヶ谷区星川1-1-1 星天qlay2階 qlaytion gallery",
   1.0, 7.1, 11.3, 0.35, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
output_path = "/home/user/test-labs/YADOKARI_Trailer_House_Proposal_v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
